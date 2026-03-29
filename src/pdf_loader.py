from pathlib import Path
import pdfplumber


def load_pdf(path: str | Path) -> list[dict]:
    """
    Extract text from a PDF file, one dict per page.

    Returns:
        List of {text, page_number, source} dicts.
        Pages with no extractable text are skipped.
    """
    path = Path(path)
    pages = []
    with pdfplumber.open(path) as pdf:
        for page_num, page in enumerate(pdf.pages, start=1):
            text = page.extract_text(x_tolerance=2, y_tolerance=2)
            if text and text.strip():
                pages.append({
                    "text": text,
                    "page_number": page_num,
                    "source": path.name,
                })
    return pages


def load_text(path: str | Path) -> list[dict]:
    """
    Load a plain text file as a single page dict.

    Returns:
        List with one entry: {text, page_number=None, source}
    """
    path = Path(path)
    text = path.read_text(encoding="utf-8", errors="replace")
    return [{"text": text, "page_number": None, "source": path.name}]


def load_docx(path: str | Path) -> list[dict]:
    """
    Extract text from a .docx file, one dict per paragraph group (treated as a page).

    Returns:
        List of {text, page_number, source} dicts.
    """
    from docx import Document
    path = Path(path)
    doc = Document(str(path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    if not paragraphs:
        return []
    # Group paragraphs into ~50-paragraph pages to preserve locality
    page_size = 50
    pages = []
    for i in range(0, len(paragraphs), page_size):
        chunk = "\n".join(paragraphs[i:i + page_size])
        pages.append({
            "text": chunk,
            "page_number": (i // page_size) + 1,
            "source": path.name,
        })
    return pages


def load_doc(path: str | Path) -> list[dict]:
    """
    Extract text from a legacy .doc file using python-docx.
    Falls back gracefully if the file is not readable.

    Returns:
        List of {text, page_number, source} dicts.
    """
    # python-docx can open many .doc files saved in newer formats
    try:
        return load_docx(path)
    except Exception:
        return []


def load_pptx(path: str | Path) -> list[dict]:
    """
    Extract text from a .pptx file, one dict per slide.

    Returns:
        List of {text, page_number, source} dicts.
        Slides with no text are skipped.
    """
    from pptx import Presentation
    path = Path(path)
    prs = Presentation(str(path))
    pages = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
        if texts:
            pages.append({
                "text": "\n".join(texts),
                "page_number": slide_num,
                "source": path.name,
            })
    return pages


def load_image(path: str | Path) -> dict:
    """
    Load an image file and return its raw bytes + metadata.

    Returns:
        A single dict with {bytes, mime_type, source, page_number, modality}
        instead of a list — images are handled separately via embed_image().
    """
    import mimetypes
    path = Path(path)
    mime_type, _ = mimetypes.guess_type(str(path))
    if mime_type is None:
        mime_type = "image/jpeg"
    return {
        "bytes": path.read_bytes(),
        "mime_type": mime_type,
        "source": path.name,
        "page_number": 0,
        "modality": "image",
    }
